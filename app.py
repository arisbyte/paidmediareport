import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import io
import anthropic

# Función para preguntar a Claude sobre los datos
def ask_claude_about_data(question, df):
    try:
        # Usar secrets si está disponible, sino pedir al usuario
        if "ANTHROPIC_API_KEY" in st.secrets:
            anthropic_api_key = st.secrets["ANTHROPIC_API_KEY"]
        else:
            return "⚠️ API key not configured. Please set ANTHROPIC_API_KEY in Streamlit secrets."
        
        client = anthropic.Anthropic(api_key=anthropic_api_key)
        
        # Preparar resumen de los datos para Claude
        data_info = {
            "rows": len(df),
            "columns": list(df.columns),
            "numeric_columns": list(df.select_dtypes(include=['number']).columns),
            "sample_data": df.head(5).to_string(),
            "summary_stats": df.describe().to_string() if len(df.select_dtypes(include=['number']).columns) > 0 else "No numeric data"
        }
        
        # Si hay columnas específicas de paid media, agregar insights
        if 'Platform' in df.columns:
            data_info["platforms"] = df['Platform'].value_counts().to_string()
        
        if 'Campaign' in df.columns:
            data_info["campaigns_count"] = df['Campaign'].nunique()
            
        message = client.messages.create(
            model="claude-3-5-sonnet-20241022",
            max_tokens=1000,
            messages=[{
                "role": "user", 
                "content": f"""You are a paid media analyst. Analyze this dataset and answer the user's question with actionable insights.

DATASET INFORMATION:
- Total rows: {data_info['rows']}
- Columns available: {data_info['columns']}
- Numeric metrics: {data_info['numeric_columns']}

SAMPLE DATA (first 5 rows):
{data_info['sample_data']}

SUMMARY STATISTICS:
{data_info['summary_stats']}

{f"PLATFORMS BREAKDOWN: {data_info['platforms']}" if 'platforms' in data_info else ""}
{f"TOTAL CAMPAIGNS: {data_info['campaigns_count']}" if 'campaigns_count' in data_info else ""}

USER QUESTION: {question}

Please provide:
1. A direct answer to the question
2. Key insights from the data
3. Actionable recommendations if applicable

Keep the response concise but informative."""
            }]
        )
        
        return message.content[0].text
        
    except Exception as e:
        return f"❌ Error calling Claude API: {str(e)}"
def create_ppt_report(df, selected_column):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    prs = Presentation()
    
    # === SLIDE 1: Portada profesional ===
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "📊 PAID MEDIA REPORT"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 56, 100)
    
    subtitle.text = f"Deep Analysis of {selected_column}\nGenerated automatically from your data"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    # === SLIDE 2: Dashboard de métricas clave ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout en blanco
    
    # Título del slide
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "🎯 KEY PERFORMANCE METRICS"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 56, 100)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Crear cajas de métricas estilo dashboard
    metrics = [
        ("TOTAL", f"{df[selected_column].sum():,.0f}", "💰"),
        ("AVERAGE", f"{df[selected_column].mean():,.0f}", "📊"),
        ("MAXIMUM", f"{df[selected_column].max():,.0f}", "🚀"),
        ("RECORDS", f"{len(df):,}", "📋")
    ]
    
    x_positions = [1, 3.5, 6, 8.5]
    for i, (label, value, emoji) in enumerate(metrics):
        # Caja de métrica
        metric_box = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(2), Inches(2), Inches(2.5))
        metric_frame = metric_box.text_frame
        metric_frame.margin_left = Inches(0.1)
        metric_frame.margin_right = Inches(0.1)
        
        # Emoji
        p1 = metric_frame.paragraphs[0]
        p1.text = emoji
        p1.font.size = Pt(36)
        p1.alignment = PP_ALIGN.CENTER
        
        # Label
        p2 = metric_frame.add_paragraph()
        p2.text = label
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(100, 100, 100)
        p2.alignment = PP_ALIGN.CENTER
        
        # Value
        p3 = metric_frame.add_paragraph()
        p3.text = value
        p3.font.size = Pt(20)
        p3.font.bold = True
        p3.font.color.rgb = RGBColor(31, 56, 100)
        p3.alignment = PP_ALIGN.CENTER
    
    # === SLIDE 3: Gráfico por plataforma (si existe) ===
    if 'Platform' in df.columns:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Título
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = f"📈 {selected_column.upper()} BY PLATFORM"
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 56, 100)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Datos para el gráfico
        platform_data = df.groupby('Platform')[selected_column].sum().sort_values(ascending=False)
        
        # Crear gráfico de barras
        chart_data = CategoryChartData()
        chart_data.categories = list(platform_data.index)
        chart_data.add_series('Values', list(platform_data.values))
        
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1.5), Inches(8), Inches(5),
            chart_data
        )
        chart = chart_shape.chart
        chart.has_legend = False
        
        # Personalizar gráfico
        chart.value_axis.has_major_gridlines = True
        chart.category_axis.tick_labels.font.size = Pt(12)
        chart.value_axis.tick_labels.font.size = Pt(12)
    
    # === SLIDE 4: Top Campaigns (si existe) ===
    if 'Campaign' in df.columns:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = f"🏆 TOP 5 CAMPAIGNS - {selected_column}"
        
        top_campaigns = df.groupby('Campaign')[selected_column].sum().sort_values(ascending=False).head(5)
        
        content = slide.placeholders[1]
        campaign_text = "🥇 " + "\n🥈 ".join([f"{camp}: {value:,.0f}" for camp, value in top_campaigns.items()])
        content.text = campaign_text
        content.text_frame.paragraphs[0].font.size = Pt(16)
    
    # === SLIDE 5: Resumen ejecutivo ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "📋 EXECUTIVE SUMMARY"
    
    # Calcular insights automáticos
    total_platforms = df['Platform'].nunique() if 'Platform' in df.columns else 0
    date_range = f"{df['Date'].min()} to {df['Date'].max()}" if 'Date' in df.columns else 'N/A'
    best_platform = df.groupby('Platform')[selected_column].sum().idxmax() if 'Platform' in df.columns else 'N/A'
    
    content = slide.placeholders[1]
    content.text = f"""📊 ANALYSIS PERIOD: {date_range}
    
🎯 TOTAL {selected_column.upper()}: {df[selected_column].sum():,.0f}

🏆 TOP PERFORMING PLATFORM: {best_platform}

📈 PLATFORMS ANALYZED: {total_platforms}

💡 RECOMMENDATION: Focus budget optimization on top-performing platforms and campaigns for maximum ROI."""
    
    content.text_frame.paragraphs[0].font.size = Pt(16)
    content.text_frame.paragraphs[0].line_spacing = 1.5
    
    # Guardar en memoria
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# Configuración de la página
st.set_page_config(page_title="Paid Media Report Processor", page_icon="📊")

# Header personalizado con fondo blanco
st.markdown(f"""
<div style="
    background-color: white;
    padding: 20px;
    border-radius: 10px;
    box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    margin-bottom: 30px;
    display: flex;
    align-items: center;
">
    <img src="https://imageshack.com/i/poUft7GXp" 
         style="height: 50px; margin-right: 20px;">
    <div>
        <h1 style="margin: 0; color: #1f3864;">Paid Media Report</h1>
        <p style="margin: 5px 0 0 0; color: #666; font-size: 16px;">Upload your Excel file with paid media data for quick analysis</p>
    </div>
</div>
""", unsafe_allow_html=True)

# Widget para subir archivo
uploaded_file = st.file_uploader(
    "Choose your Excel file", 
    type=['xlsx', 'xls'],
    help="File should contain columns like: Campaign, Impressions, Clicks, Cost"
)

# Si hay archivo subido
if uploaded_file is not None:
    try:
        # Leer el Excel
        df = pd.read_excel(uploaded_file)
        
        # Información básica
        st.subheader("📋 Data Overview")
        col1, col2 = st.columns(2)
        with col1:
            st.metric("Total Rows", len(df))
        with col2:
            st.metric("Columns", len(df.columns))
        
        st.write("**Columns found:**", list(df.columns))
        
        # Mostrar muestra de datos
        st.subheader("📊 Data Sample")
        st.dataframe(df.head(10))
        
        # Análisis básico si hay columnas numéricas
        numeric_columns = df.select_dtypes(include=['number']).columns
        
        if len(numeric_columns) > 0:
            st.subheader("📈 Quick Analysis")
            
            # Seleccionar columna para analizar
            selected_column = st.selectbox("Select a metric to analyze:", numeric_columns)
            
            if selected_column:
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("Total", f"{df[selected_column].sum():,.0f}")
                with col2:
                    st.metric("Average", f"{df[selected_column].mean():,.0f}")
                with col3:
                    st.metric("Max", f"{df[selected_column].max():,.0f}")
                
                # Gráfico simple
                st.subheader(f"📊 {selected_column} Distribution")
                fig, ax = plt.subplots(figsize=(10, 6))
                ax.hist(df[selected_column].dropna(), bins=20, alpha=0.7, color='skyblue')
                ax.set_xlabel(selected_column)
                ax.set_ylabel('Frequency')
                ax.set_title(f'Distribution of {selected_column}')
                st.pyplot(fig)
        
        # === NUEVA SECCIÓN: AI INSIGHTS ===
        st.subheader("🤖 AI Data Analysis")
        
        # Verificar si hay API key configurada
        if "ANTHROPIC_API_KEY" in st.secrets:
            st.write("Ask Claude anything about your paid media data:")
            
            # Input para la pregunta
            user_question = st.text_input(
                "Your question:", 
                placeholder="e.g., Which platform has the best ROI? What trends do you see?"
            )
            
            # Botón para hacer la pregunta
            if user_question and st.button("🔍 Ask Claude", type="primary"):
                with st.spinner("Claude is analyzing your data..."):
                    response = ask_claude_about_data(user_question, df)
                    
                    # Mostrar respuesta
                    st.success("📝 Claude's Analysis:")
                    st.write(response)
                    
        else:
            st.info("💡 AI analysis available when API key is configured in Streamlit secrets.")
            
            # Mostrar ejemplos de preguntas que podrían hacer
            st.write("**Example questions you could ask:**")
            st.write("• Which platform has the best performance?")
            st.write("• What patterns do you see in the data?")
            st.write("• Which campaigns should I optimize?")
            st.write("• What's driving the highest costs?")
                
        # Descarga de datos procesados
        st.subheader("💾 Download Options")
        
        col1, col2 = st.columns(2)
        
        with col1:
            csv = df.to_csv(index=False)
            st.download_button(
                label="📄 Download as CSV",
                data=csv,
                file_name='processed_paid_media_data.csv',
                mime='text/csv'
            )
        
        with col2:
            # AQUÍ SE AGREGA EL BOTÓN DE POWERPOINT
            if len(numeric_columns) > 0 and 'selected_column' in locals():
                ppt_file = create_ppt_report(df, selected_column)
                st.download_button(
                    label="📊 Download PowerPoint Report",
                    data=ppt_file,
                    file_name=f'paid_media_report_{selected_column}.pptx',
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            
    except Exception as e:
        st.error(f"❌ Error processing file: {str(e)}")
        st.write("Make sure it's a valid Excel file")

else:
    # Mostrar ejemplo de datos esperados
    st.info("👆 Upload an Excel file to get started")
    
    st.subheader("📝 Expected File Format:")
    sample_data = {
        'Campaign': ['Google Ads - Search', 'Facebook Ads', 'Instagram Ads', 'YouTube Ads'],
        'Impressions': [15000, 25000, 18000, 12000],
        'Clicks': [850, 1200, 950, 600],
        'Cost': [420.50, 680.25, 510.75, 290.30]
    }
    sample_df = pd.DataFrame(sample_data)
    st.dataframe(sample_df)
    
    st.write("**The app will automatically detect and analyze any numeric columns in your data.**")